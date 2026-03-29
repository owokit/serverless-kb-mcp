"""
EN: Tests for document extraction from Markdown, DOCX, and PPTX formats.
CN: 娴嬭瘯浠?Markdown銆丏OCX 鍜?PPTX 鏍煎紡鎻愬彇鏂囨。銆?
"""

from io import BytesIO

from docx import Document
from pptx import Presentation
from pptx.util import Inches

from serverless_mcp.extract import extractors as extractor_module
from serverless_mcp.extract.extractors import DocumentExtractor
from serverless_mcp.extract.policy import DEFAULT_POLICY
from serverless_mcp.domain.models import S3ObjectRef


def test_markdown_extraction_keeps_section_path_and_version() -> None:
    """
    EN: Markdown extraction keeps section path and version.
    CN: 楠岃瘉 Markdown 鎻愬彇淇濈暀 section_path 鍜?version_id銆?
    """
    extractor = DocumentExtractor()
    manifest = extractor.extract(
        source=S3ObjectRef(
            tenant_id="tenant-a",
            bucket="bucket-a",
            key="docs/guide.md",
            version_id="v123",
        ),
        body=b"# Intro\nhello world\n\n## Details\nmore text\n",
        safe_text_tokens=128,
    )

    assert manifest.source.version_id == "v123"
    assert manifest.doc_type == "md"
    assert [chunk.section_path for chunk in manifest.chunks] == [("Intro",), ("Intro", "Details")]
    assert all(chunk.chunk_type == "section_text_chunk" for chunk in manifest.chunks)


def test_markdown_extraction_uses_default_gemini_safe_chunk_size() -> None:
    """
    EN: Markdown extraction uses default gemini safe chunk size.
    CN: 同上。
    """
    extractor = DocumentExtractor()
    manifest = extractor.extract(
        source=S3ObjectRef(
            tenant_id="tenant-a",
            bucket="bucket-a",
            key="docs/long.md",
            version_id="v123",
        ),
        body=b"# Intro\n" + ("word " * 2200).encode("utf-8"),
    )

    assert len(manifest.chunks) > 1
    assert all(chunk.token_estimate <= DEFAULT_POLICY.safe_text_tokens for chunk in manifest.chunks)


def test_docx_extraction_uses_section_headings_and_tables() -> None:
    """
    EN: Docx extraction uses section headings and tables.
    CN: 同上。
    """
    buffer = BytesIO()
    document = Document()
    document.add_heading("Intro", level=1)
    document.add_paragraph("hello world")
    document.add_heading("Details", level=2)
    document.add_paragraph("more text")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "left"
    table.cell(0, 1).text = "right"
    document.save(buffer)

    extractor = DocumentExtractor()
    manifest = extractor.extract(
        source=S3ObjectRef(
            tenant_id="tenant-a",
            bucket="bucket-a",
            key="docs/guide.docx",
            version_id="v123",
        ),
        body=buffer.getvalue(),
        safe_text_tokens=128,
    )

    assert manifest.doc_type == "docx"
    assert manifest.metadata["source_format"] == "python-docx"
    assert [chunk.section_path for chunk in manifest.chunks] == [("Intro",), ("Intro", "Details")]
    assert manifest.chunks[0].text.startswith("# Intro")
    assert "hello world" in manifest.chunks[0].text
    assert "left" in "\n".join(chunk.text for chunk in manifest.chunks)
    assert "right" in "\n".join(chunk.text for chunk in manifest.chunks)


def test_pptx_extraction_uses_slide_segments() -> None:
    """
    EN: Pptx extraction uses slide segments.
    CN: 同上。
    """
    buffer = BytesIO()
    presentation = Presentation()
    slide1 = presentation.slides.add_slide(presentation.slide_layouts[5])
    textbox = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    textbox.text = "Title"
    slide2 = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide2.shapes.title.text = "Second"
    slide2.placeholders[1].text = "Bullet one\nBullet two"
    presentation.save(buffer)

    extractor = DocumentExtractor()
    manifest = extractor.extract(
        source=S3ObjectRef(
            tenant_id="tenant-a",
            bucket="bucket-a",
            key="docs/deck.pptx",
            version_id="v123",
        ),
        body=buffer.getvalue(),
        safe_text_tokens=128,
    )

    assert manifest.doc_type == "pptx"
    assert manifest.metadata["source_format"] == "python-pptx"
    assert [chunk.section_path for chunk in manifest.chunks] == [("slide-1",), ("slide-2",)]
    assert manifest.chunks[0].text == "Title"
    assert "Second" in manifest.chunks[1].text


def test_pptx_extraction_parses_presentation_once(monkeypatch) -> None:
    """
    EN: PPTX extraction should only parse the presentation once.
    CN: PPTX 抽取应当只解析 presentation 一次。
    """
    buffer = BytesIO()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    textbox.text = "Title"
    presentation.save(buffer)

    calls = {"count": 0}
    real_presentation = Presentation

    def _counting_presentation(*args, **kwargs):
        calls["count"] += 1
        return real_presentation(*args, **kwargs)

    monkeypatch.setattr(extractor_module, "Presentation", _counting_presentation)

    extractor = DocumentExtractor()
    extractor.extract(
        source=S3ObjectRef(
            tenant_id="tenant-a",
            bucket="bucket-a",
            key="docs/deck.pptx",
            version_id="v123",
        ),
        body=buffer.getvalue(),
        safe_text_tokens=128,
    )

    assert calls["count"] == 1


def test_pdf_extraction_uses_configured_page_part_windows(monkeypatch) -> None:
    """
    EN: PDF extraction should use the configured page-part window size when building window chunks.
    CN: PDF 抽取应当在构造 window chunk 时使用配置的页分片大小。
    """
    class _Page:
        def __init__(self, text: str) -> None:
            self._text = text

        def extract_text(self) -> str:
            return self._text

    class _FakePdfReader:
        def __init__(self, _buffer) -> None:
            self.pages = [_Page("page one"), _Page("page two"), _Page("page three")]

    monkeypatch.setattr(extractor_module, "PdfReader", _FakePdfReader)
    monkeypatch.setattr(extractor_module, "estimate_tokens", lambda text: 80 if "\n\n" not in text else 100)

    extractor = DocumentExtractor()
    manifest = extractor.extract(
        source=S3ObjectRef(
            tenant_id="tenant-a",
            bucket="bucket-a",
            key="docs/guide.pdf",
            version_id="v123",
        ),
        body=b"%PDF-1.7",
        safe_text_tokens=200,
        max_pdf_pages_per_part=2,
    )

    assert any(chunk.metadata.get("window_reason") == "page_part_window" for chunk in manifest.chunks)
    assert any(chunk.page_span == (1, 2) for chunk in manifest.chunks if chunk.chunk_type == "window_pdf_chunk")
