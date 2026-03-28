"""
EN: Document extractors for multiple formats including PDF, DOCX, PPTX, and Markdown.
CN: 同上。
"""
from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
import re
from typing import Protocol

from docx import Document as load_docx_document
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table
from docx.text.paragraph import Paragraph
from pypdf import PdfReader
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from serverless_mcp.domain.format_specs import get_format_spec
from serverless_mcp.domain.models import ChunkManifest, ExtractedAsset, ExtractedChunk, S3ObjectRef
from serverless_mcp.extract.policy import (
    DEFAULT_POLICY,
    estimate_tokens,
    expand_oversized_chunks,
    normalize_text,
    section_hint_from_markdown,
)


# EN: Regex to match PPTX slide number markers emitted by the Markdown converter.
# CN: 同上。
PPTX_SLIDE_MARKER_PATTERN = re.compile(r"^<!-- Slide number: (\d+) -->$")


class OCRClient(Protocol):
    """EN: Protocol for OCR client implementations.
    CN: 同上。
    """

    def extract_text(self, *, mime_type: str, payload: bytes) -> str: ...


class PageRasterizer(Protocol):
    """EN: Protocol for PDF page rendering implementations.
    CN: 同上。
    """

    def render_pdf_page(self, *, pdf_bytes: bytes, page_number: int) -> bytes | None: ...


@dataclass(slots=True)
class ExtractionContext:
    """EN: Context container for extraction operations with optional OCR and rasterization support.
    CN: 同上。
    """

    source: S3ObjectRef
    safe_text_tokens: int
    ocr_client: OCRClient | None = None
    page_rasterizer: PageRasterizer | None = None


class UnsupportedDocumentTypeError(ValueError):
    """EN: Raised when document format is not supported by any extractor.
    CN: 同上。
    """


class DocumentExtractor:
    """EN: Multi-format document extractor supporting MD, DOCX, PPTX, and PDF with chunking.
    CN: 同上。
    """

    def extract(
        self,
        *,
        source: S3ObjectRef,
        body: bytes,
        safe_text_tokens: int = DEFAULT_POLICY.safe_text_tokens,
        max_pdf_pages_per_part: int = DEFAULT_POLICY.max_pdf_pages_per_part,
    ) -> ChunkManifest:
        """
        EN: Extract chunks from document based on file extension and apply oversized chunk expansion.
    CN: 根据文件扩展名提取 chunk，并对超大 chunk 做扩展拆分。

        Args:
            source:
                EN: S3 object reference identifying the document by bucket, key, and version_id.
                CN: 通过 bucket、key 和 version_id 标识文档的 S3 对象引用。
            body:
                EN: Raw document bytes fetched from S3.
                CN: 从 S3 获取的原始文档字节。
            safe_text_tokens:
                EN: Maximum safe token count per chunk; oversized chunks are split further.
                CN: 同上。

        Returns:
            EN: Chunk manifest with extracted text chunks and optional assets.
            CN: 包含已提取文本 chunk 和可选资产的 chunk manifest。

        Raises:
            UnsupportedDocumentTypeError:
                EN: When the file extension is not recognized or is a legacy Office binary format.
                CN: 当文件扩展名无法识别或属于旧式 Office 二进制格式时抛出。
        """

        # EN: Build extraction context; dispatched by extension to format-specific handlers.
        # CN: 构建提取上下文；随后按扩展名分发到具体格式处理器。
        ctx = ExtractionContext(source=source, safe_text_tokens=safe_text_tokens)
        extension = source.extension

        if extension == "md":
            manifest = self._extract_markdown(ctx, body)
        elif extension == "docx":
            manifest = self._extract_docx(ctx, body)
        elif extension == "pptx":
            manifest = self._extract_pptx(ctx, body)
        elif extension == "pdf":
            manifest = self._extract_pdf(ctx, body, max_pdf_pages_per_part=max_pdf_pages_per_part)
        elif extension in {"doc", "ppt"}:
            raise UnsupportedDocumentTypeError(
                f"Legacy Office binary format '.{extension}' is not supported directly. Convert to OOXML first."
            )
        else:
            raise UnsupportedDocumentTypeError(f"Unsupported extension: {extension or '<none>'}")

        # EN: Post-processing pass 鈥?split oversized chunks to respect embedding token limits.
        # CN: 同上。
        manifest.chunks = expand_oversized_chunks(manifest.chunks, safe_text_tokens=safe_text_tokens)
        return manifest

    def _extract_markdown(self, ctx: ExtractionContext, body: bytes) -> ChunkManifest:
        """
        EN: Extract Markdown content by parsing heading structure into section-based chunks.
        CN: 同上。

        Args:
            ctx:
                EN: Extraction context carrying the source reference and token limits.
                CN: 携带源引用和 token 限制的提取上下文。
            body:
                EN: Raw Markdown document bytes.
                CN: 同上。

        Returns:
            EN: Chunk manifest with section_text_chunk entries per heading section.
            CN: 同上。
        """
        spec = get_format_spec(doc_type="md", source_format="markdown")
        text = body.decode("utf-8")
        sections = section_hint_from_markdown(text)
        chunks: list[ExtractedChunk] = []

        if not sections:
            # EN: Fallback to a single normalized chunk when no headings are detected.
            # CN: 同上。
            normalized = normalize_text(text)
            sections = [((), normalized)] if normalized else []

        for index, (path, section_text) in enumerate(sections, start=1):
            chunks.append(
                ExtractedChunk(
                    chunk_id=f"chunk#{index:06d}",
                    chunk_type="section_text_chunk",
                    text=section_text,
                    doc_type="md",
                    token_estimate=estimate_tokens(section_text),
                    section_path=path,
                    metadata=spec.chunk_metadata(),
                )
            )

        return ChunkManifest(
            source=ctx.source,
            doc_type="md",
            chunks=chunks,
            metadata=spec.manifest_metadata(section_count=len(chunks)),
        )

    def _extract_docx(self, ctx: ExtractionContext, body: bytes) -> ChunkManifest:
        """
        EN: Extract DOCX content by converting to Markdown first, then splitting by heading sections.
        CN: 同上。

        Args:
            ctx:
                EN: Extraction context carrying the source reference and token limits.
                CN: 携带源引用和 token 限制的提取上下文。
            body:
                EN: Raw DOCX document bytes.
                CN: 同上。

        Returns:
            EN: Chunk manifest with section_text_chunk entries per heading section.
            CN: 同上。
        """
        spec = get_format_spec(doc_type="docx", source_format="python-docx")
        markdown_text = _convert_docx_to_markdown(body)
        sections = section_hint_from_markdown(markdown_text)
        chunks: list[ExtractedChunk] = []

        if not sections:
            # EN: Fallback to a single normalized chunk when no headings are detected.
            # CN: 同上。
            normalized = normalize_text(markdown_text)
            sections = [((), normalized)] if normalized else []

        for index, (path, section_text) in enumerate(sections, start=1):
            chunks.append(
                ExtractedChunk(
                    chunk_id=f"chunk#{index:06d}",
                    chunk_type="section_text_chunk",
                    text=section_text,
                    doc_type="docx",
                    token_estimate=estimate_tokens(section_text),
                    section_path=path,
                    metadata=spec.chunk_metadata(),
                )
            )

        return ChunkManifest(
            source=ctx.source,
            doc_type="docx",
            chunks=chunks,
            metadata=spec.manifest_metadata(section_count=len(chunks)),
        )

    def _extract_pptx(self, ctx: ExtractionContext, body: bytes) -> ChunkManifest:
        """
        EN: Extract PPTX slides as both text chunks and image assets, including speaker notes.
        CN: 同上。

        Args:
            ctx:
                EN: Extraction context carrying the source reference and token limits.
                CN: 携带源引用和 token 限制的提取上下文。
            body:
                EN: Raw PPTX document bytes.
                CN: 同上。

        Returns:
            EN: Chunk manifest with slide_text_chunk and slide_image_chunk entries.
            CN: 包含 slide_text_chunk 和 slide_image_chunk 条目的 chunk manifest。
        """
        spec = get_format_spec(doc_type="pptx", source_format="python-pptx")
        presentation = Presentation(BytesIO(body))
        markdown_text = _convert_pptx_to_markdown(presentation)
        slide_markdowns = _split_pptx_markdown(markdown_text)
        chunks: list[ExtractedChunk] = []
        assets: list[ExtractedAsset] = []
        asset_index = 0

        # EN: First pass 鈥?extract embedded images as assets per slide.
        # CN: 第一遍，按幻灯片提取嵌入图片作为资产。
        for slide_no, slide in enumerate(presentation.slides, start=1):
            image_count = 0

            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_count += 1
                    asset_index += 1
                    image = shape.image
                    assets.append(
                        ExtractedAsset(
                            asset_id=f"asset#{asset_index:06d}",
                            chunk_type="slide_image_chunk",
                            mime_type=image.content_type,
                            payload=image.blob,
                            slide_no=slide_no,
                            metadata=spec.asset_metadata(
                                filename=image.filename,
                                image_count_on_slide=image_count,
                            ),
                        )
                    )

            # EN: Merge slide text with speaker notes into a single chunk per slide.
            # CN: 同上。
            slide_text = slide_markdowns.get(slide_no, "")
            if slide.has_notes_slide and slide.notes_slide:
                notes = [
                    shape.text.strip()
                    for shape in slide.notes_slide.shapes
                    if hasattr(shape, "text") and shape.text.strip()
                ]
                notes_text = normalize_text("\n".join(notes))
                if notes_text:
                    slide_text = normalize_text("\n".join(part for part in [slide_text, notes_text] if part))

            chunks.append(
                ExtractedChunk(
                    chunk_id=f"chunk#{slide_no:06d}",
                    chunk_type="slide_text_chunk",
                    text=slide_text,
                    doc_type="pptx",
                    token_estimate=estimate_tokens(slide_text),
                    slide_no=slide_no,
                    section_path=(f"slide-{slide_no}",),
                    metadata=spec.chunk_metadata(image_count=image_count, has_notes=bool(slide.has_notes_slide and slide.notes_slide)),
                )
            )

        return ChunkManifest(
            source=ctx.source,
            doc_type="pptx",
            chunks=chunks,
            assets=assets,
            metadata=spec.manifest_metadata(slide_count=len(chunks), image_asset_count=len(assets)),
        )

    def _extract_pdf(
        self,
        ctx: ExtractionContext,
        body: bytes,
        *,
        max_pdf_pages_per_part: int,
    ) -> ChunkManifest:
        """
        EN: Extract PDF pages as text chunks with optional page-image assets and short-page window merging.
        CN: 同上。

        Args:
            ctx:
                EN: Extraction context carrying the source reference, token limits, and optional rasterizer.
                CN: 携带源引用、token 限制和可选光栅化器的提取上下文。
            body:
                EN: Raw PDF document bytes.
                CN: 同上。

        Returns:
            EN: Chunk manifest with page_text_chunk, window_pdf_chunk, and page_image_chunk assets.
            CN: 包含 page_text_chunk、window_pdf_chunk 和 page_image_chunk 资产的 chunk manifest。
        """
        spec = get_format_spec(doc_type="pdf", source_format="pdf")
        reader = PdfReader(BytesIO(body))
        page_chunks: list[ExtractedChunk] = []
        window_chunks: list[ExtractedChunk] = []
        assets: list[ExtractedAsset] = []
        chunk_index = 0
        visual_pages: list[int] = []
        max_pdf_pages_per_part = max(1, max_pdf_pages_per_part)

        for page_no, page in enumerate(reader.pages, start=1):
            text = normalize_text(page.extract_text() or "")
            # EN: Pages with extractable text become page_text_chunk entries.
            # CN: 同上。
            if text:
                chunk_index += 1
                page_chunks.append(
                    ExtractedChunk(
                        chunk_id=f"chunk#{chunk_index:06d}",
                        chunk_type="page_text_chunk",
                        text=text,
                        doc_type="pdf",
                        token_estimate=estimate_tokens(text),
                        page_no=page_no,
                        page_span=(page_no, page_no),
                        section_path=(f"page-{page_no}",),
                        metadata=spec.chunk_metadata(),
                    )
                )
            else:
                # EN: Track pages with no extractable text as visual-only pages.
                # CN: 同上。
                visual_pages.append(page_no)

            # EN: Optionally render page images for multimodal embedding when a rasterizer is available.
            # CN: 同上。
            if ctx.page_rasterizer is not None:
                rendered = ctx.page_rasterizer.render_pdf_page(pdf_bytes=body, page_number=page_no)
                if rendered:
                    assets.append(
                        ExtractedAsset(
                            asset_id=f"asset#{page_no:06d}",
                            chunk_type="page_image_chunk",
                            mime_type="image/png",
                            payload=rendered,
                            page_no=page_no,
                            metadata=spec.asset_metadata(),
                        )
                    )

        # EN: Window merge 鈥?combine adjacent short-text pages into cross-page chunks.
        # CN: 窗口合并，将相邻短文本页合并成跨页 chunk。
        for current, nxt in zip(page_chunks, page_chunks[1:]):
            if current.page_no is None or nxt.page_no is None:
                continue
            combined = normalize_text(f"{current.text}\n\n{nxt.text}")
            if estimate_tokens(combined) > ctx.safe_text_tokens:
                continue
            if current.token_estimate < 120 or nxt.token_estimate < 120:
                chunk_index += 1
                window_chunks.append(
                    ExtractedChunk(
                        chunk_id=f"chunk#{chunk_index:06d}",
                        chunk_type="window_pdf_chunk",
                        text=combined,
                        doc_type="pdf",
                        token_estimate=estimate_tokens(combined),
                        page_span=(current.page_no, nxt.page_no),
                        section_path=(f"pages-{current.page_no}-{nxt.page_no}",),
                        metadata=spec.chunk_metadata(window_reason="short_adjacent_pages"),
                    )
                )

        # EN: Larger PDF parts use the configured page cap to build coarser part windows when they still fit.
        # CN: 同上。
        for start_index in range(0, len(page_chunks), max_pdf_pages_per_part):
            part = page_chunks[start_index : start_index + max_pdf_pages_per_part]
            if len(part) < 2:
                continue
            combined = normalize_text("\n\n".join(chunk.text for chunk in part))
            if estimate_tokens(combined) > ctx.safe_text_tokens:
                continue
            first_page = part[0].page_no
            last_page = part[-1].page_no
            if first_page is None or last_page is None:
                continue
            window_chunks.append(
                ExtractedChunk(
                    chunk_id=f"chunk#{chunk_index + 1:06d}",
                    chunk_type="window_pdf_chunk",
                    text=combined,
                    doc_type="pdf",
                    token_estimate=estimate_tokens(combined),
                    page_span=(first_page, last_page),
                    section_path=(f"pages-{first_page}-{last_page}",),
                    metadata=spec.chunk_metadata(window_reason="page_part_window"),
                )
            )
            chunk_index += 1

        return ChunkManifest(
            source=ctx.source,
            doc_type="pdf",
            chunks=page_chunks + window_chunks,
            assets=assets,
            metadata={
                "page_count": len(reader.pages),
                "visual_page_numbers": visual_pages,
                "page_image_asset_count": len(assets),
                "source_format": spec.source_format,
            },
        )


def _convert_docx_to_markdown(body: bytes) -> str:
    """
    EN: Convert DOCX body bytes to a Markdown string using python-docx structural iteration.
    CN: 同上。

    Args:
        body:
            EN: Raw DOCX document bytes.
            CN: 同上。

    Returns:
        EN: Normalized Markdown text preserving heading levels and table formatting.
        CN: 同上。
    """
    document = load_docx_document(BytesIO(body))
    lines: list[str] = []
    # EN: Iterate document body blocks in document order (paragraphs and tables).
    # CN: 同上。
    for block in _iter_docx_blocks(document):
        if isinstance(block, Paragraph):
            rendered = _paragraph_to_markdown(block)
            if rendered:
                lines.append(rendered)
            continue
        if isinstance(block, Table):
            lines.extend(_table_to_markdown(block))
    return normalize_text("\n".join(lines))


def _iter_docx_blocks(document) -> list[Paragraph | Table]:
    """
    EN: Iterate over top-level body elements preserving document order of paragraphs and tables.
    CN: 同上。

    Args:
        document:
            EN: Loaded python-docx Document object.
            CN: 已加载的 python-docx Document 对象。

    Returns:
        EN: Ordered list of Paragraph and Table blocks.
        CN: 按顺序排列的 Paragraph 和 Table 块列表。
    """
    blocks: list[Paragraph | Table] = []
    # EN: Classify each XML child element by its underlying type (CT_P for paragraphs, CT_Tbl for tables).
    # CN: 同上。
    for child in document.element.body.iterchildren():
        if isinstance(child, CT_P):
            blocks.append(Paragraph(child, document))
        elif isinstance(child, CT_Tbl):
            blocks.append(Table(child, document))
    return blocks


def _paragraph_to_markdown(paragraph: Paragraph) -> str:
    """
    EN: Convert a single DOCX paragraph to Markdown, mapping heading styles and list styles.
    CN: 同上。

    Args:
        paragraph:
            EN: A python-docx Paragraph instance.
            CN: 一个 python-docx Paragraph 实例。

    Returns:
        EN: Markdown-formatted text line, or empty string if the paragraph has no content.
        CN: 同上。
    """
    text = normalize_text(paragraph.text)
    if not text:
        return ""

    style_name = (paragraph.style.name if paragraph.style is not None else "").strip()
    heading_level = _heading_level_from_style_name(style_name)
    # EN: Render heading prefix (# / ## / etc.) when the paragraph style is a recognized heading.
    # CN: 同上。
    if heading_level:
        return f"{'#' * heading_level} {text}"

    if style_name.lower().startswith("list"):
        return f"- {text}"

    return text


def _heading_level_from_style_name(style_name: str) -> int:
    """
    EN: Map a Word style name to a Markdown heading level (1-6), returning 0 for non-heading styles.
    CN: 同上。

    Args:
        style_name:
            EN: Style name from a Word paragraph (e.g., "Heading 1", "Title").
            CN: 同上。

    Returns:
        EN: Heading level integer from 1 to 6, or 0 if not a heading style.
        CN: 1 到 6 的标题级别整数；若不是标题样式则返回 0。
    """
    if style_name.lower() == "title":
        return 1

    match = re.match(r"Heading\s+(\d+)$", style_name, flags=re.IGNORECASE)
    if not match:
        return 0

    return max(1, min(6, int(match.group(1))))


def _table_to_markdown(table: Table) -> list[str]:
    """
    EN: Convert a DOCX table to Markdown pipe-delimited lines with a header separator row.
    CN: 同上。

    Args:
        table:
            EN: A python-docx Table instance.
            CN: 一个 python-docx Table 实例。

    Returns:
        EN: List of Markdown table lines including header, separator, and data rows.
        CN: 同上。
    """
    rows: list[list[str]] = []
    # EN: Normalize cell text and skip entirely empty rows.
    # CN: 同上。
    for row in table.rows:
        cells = [normalize_text(cell.text) for cell in row.cells]
        if any(cells):
            rows.append(cells)

    if not rows:
        return []

    width = max(len(row) for row in rows)
    normalized_rows = [row + [""] * (width - len(row)) for row in rows]
    lines = ["| " + " | ".join(normalized_rows[0]) + " |"]
    lines.append("| " + " | ".join("---" for _ in range(width)) + " |")
    for row in normalized_rows[1:]:
        lines.append("| " + " | ".join(row) + " |")
    return lines


def _convert_pptx_to_markdown(presentation) -> str:
    """
    EN: Convert PPTX body bytes to Markdown with per-slide markers for later splitting.
    CN: 同上。

    Args:
        presentation:
            EN: Loaded python-pptx Presentation instance.
            CN: 已加载的 python-pptx Presentation 实例。

    Returns:
        EN: Normalized Markdown text with <!-- Slide number: N --> markers between slides.
        CN: 同上。
    """
    slide_blocks: list[str] = []
    # EN: Each slide block starts with a <!-- Slide number: N --> marker for downstream indexing.
    # CN: 同上。
    for slide_no, slide in enumerate(presentation.slides, start=1):
        lines = [f"<!-- Slide number: {slide_no} -->"]
        lines.extend(_slide_text_lines(slide))
        slide_blocks.append("\n".join(line for line in lines if line.strip()))
    return normalize_text("\n\n".join(slide_blocks))


def _slide_text_lines(slide) -> list[str]:
    """
    EN: Extract all text lines from a slide's shapes and speaker notes.
    CN: 同上。

    Args:
        slide:
            EN: A python-pptx Slide object.
            CN: 一个 python-pptx Slide 对象。

    Returns:
        EN: List of non-empty text lines from text frames and notes slide.
        CN: 同上。
    """
    lines: list[str] = []
    # EN: Collect text from shapes that have a text_frame attribute.
    # CN: 同上。
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = normalize_text(getattr(shape.text_frame, "text", "") or "")
        if text:
            lines.extend(line for line in text.splitlines() if line.strip())

    # EN: Append speaker notes text if present on the slide.
    # CN: 同上。
    if slide.has_notes_slide and slide.notes_slide:
        for shape in slide.notes_slide.shapes:
            text = normalize_text(getattr(shape, "text", "") or "")
            if text:
                lines.extend(line for line in text.splitlines() if line.strip())

    return lines


def _split_pptx_markdown(markdown_text: str) -> dict[int, str]:
    """EN: Split PowerPoint Markdown into slide-indexed fragments.
    CN: 同上。
    """

    segments: dict[int, list[str]] = {}
    slide_no: int | None = None

    for line in normalize_text(markdown_text).splitlines():
        marker = PPTX_SLIDE_MARKER_PATTERN.match(line.strip())
        if marker is not None:
            slide_no = int(marker.group(1))
            segments.setdefault(slide_no, [])
            continue
        if slide_no is None:
            continue
        segments.setdefault(slide_no, []).append(line)

    return {key: normalize_text("\n".join(lines)) for key, lines in segments.items()}
